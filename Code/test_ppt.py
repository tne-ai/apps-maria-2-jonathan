import json
from io import BytesIO
from pptx import Presentation 

def ppt_maker(content, file_name):
    # create actual powerpoint object
    ppt = Presentation()

    #Iterate through each section
    for section in content["sections"]:
        json_title = section["title"]
        json_content = section["content"]
        json_content_type = section["type"]

        if json_content_type == "raw text":
            #have text slide
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            slide_title = slide.shapes.title
            slide_text = slide.shapes.placeholders[1]

            slide_title.text = json_title
            slide_text.text = json_content

    # save powerpoint
    ppt.save(file_name)

with open('ppt_test1.txt', 'r') as file:
    txt_contents = file.read()

#turn it into dictionary
json_contents = json.loads(txt_contents)

file_name = json_contents["document_filename"]
print(file_name)
content = json_contents

#call pdf_maker method
file = ppt_maker(content, file_name)

